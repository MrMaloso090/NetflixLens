def Data_presentation():
    global PAGE

    import os
    import sys
    import sqlite3
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt
    import requests
    from PIL import Image



    # SE TOMA LA CONEXION CON LA UBICACION DEL ARCHIVO.
    # THE CONNECTION WITH THE FILE LOCATION IS TAKEN.
    try:
        path_created = sqlite3.connect('Path_File.sqlite')
        cur_path = path_created.cursor()

        cur_path.execute('''
            SELECT path FROM path_data WHERE id = (?)''', (1, ))

        PATH = cur_path.fetchone()[0]
        if len(PATH) < 20:
            raise
        print(PATH)
        path_created.close()
    except:
        print('==============================================================\n')
        print('======================== ERROR ===============================')
        print('==================== PATH NOT FOUND ==========================')
        print('================= DO  ALL STEPS IN ORDER =====================\n')
        print('==============================================================\n')
        return



    # VERIFICACION DE LA EXISTENCIA DE LOS GRAFICOS.
    # VERIFICATION OF THE EXISTENCE OF THE GRAPHS.
    if not (os.path.exists(PATH + 'Global_Graphs')) and (os.path.exists(PATH + 'Tv_Graphs')) and (os.path.exists(PATH + 'Movie_Graphs')):
        print('==============================================================\n')
        print('======================== ERROR ===============================')
        print('================ GRAPHS FORLDERS NOT FOUND ===================')
        print('================= DO  ALL STEPS IN ORDER =====================\n')
        print('==============================================================\n')
        return


    # ASIGNACION DE LOS CAMINOS A LAS CARPETAS A SU RESPECTIVA VARIABLE.
    # ASSIGNMENT OF THE FOLDER PATHS TO THEIR RESPECTIVE VARIABLE.
    FOLDER_GLOBAL = (PATH + 'Global_Graphs//')
    FOLDER_TV = (PATH + 'Tv_Graphs//')
    FOLDER_MOVIE = (PATH + 'Movie_Graphs//')
    FOLDER_PRESENTATION = '_data//Presentation//'



    ############################## PRESENTATION.
    ############################## PRESENTATION.
    diapo_ = Presentation()
    diapo_.slide_width = Inches(16)
    diapo_.slide_height = Inches(9)


    # PAGINA 1. TITULO
    # PAGE 1. TITLE
    page_1 = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    page_1.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_1.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)


    gm_1_1 = page_1.shapes.add_textbox(Inches(3.5) , Inches(2.8) , Inches(6) , Inches(1))
    p1_t1 = gm_1_1.text_frame.paragraphs[0]
    p1_t1.text = 'Trend Analysis of Netflix Content'
    p1_t1.font.name = 'Impact'
    p1_t1.font.size = Pt(48)
    p1_t1.font.color.rgb = RGBColor(211, 211, 211)

    gm_1_2 = page_1.shapes.add_textbox(Inches(5) , Inches(4) , Inches(6) , Inches(1))
    p1_t2 = gm_1_2.text_frame.paragraphs[0]
    p1_t2.text = 'Exploration of the most and least viewed, popularity by genres and countries, \nand the relationship between views and playback time.'
    p1_t2.font.name = 'Arial Narrow'
    p1_t2.font.size = Pt(32)
    p1_t2.font.color.rgb = RGBColor(211, 211, 211)
    p1_t2.alignment = PP_ALIGN.CENTER


    gm_f = page_1.shapes.add_textbox(Inches(0.5) , Inches(0.1) , Inches(6) , Inches(1))
    p1_tf = gm_f.text_frame.paragraphs[0]

    p1_tf.text = '''
    __________________________________________________

    Code Author: Daniel Sanchez Velasquez - TakuSan.
    Email: daniel.sanchez.velasquez090@gmail.com
    __________________________________________________
    '''

    p1_tf.font.name = 'Arial Narrow'
    p1_tf.font.size = Pt(18)
    p1_tf.font.color.rgb = RGBColor(211, 211, 211)
    p1_tf.alignment = PP_ALIGN.LEFT


    page_1.shapes.add_picture(FOLDER_PRESENTATION + 'TMDB.png' , left= Inches(13.5) , top= Inches(6.5) , width = Inches(2) , height = Inches(2))

    page_1.shapes.add_picture(FOLDER_PRESENTATION + 'NETFLIX.png' , left= Inches(0.5) , top= Inches(6.5) , width = Inches(1) , height = Inches(2))



    # PAGINA 1.5 TITULO
    # PAGE 1.5 TITLE
    page_p = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    page_p.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_1.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)


    gm_p = page_p.shapes.add_textbox(Inches(0.5) , Inches(0.1) , Inches(6) , Inches(1))
    p1_tp = gm_p.text_frame.paragraphs[0]

    p1_tp.text = '''
    This project offers a comprehensive, automated solution for analyzing Netflix's vast catalog of audiovisual content. Leveraging the Netflix Data Dump, introduced in the 
    second half of 2023, the program provides a thorough exploration of titles, viewing metrics, and additional metadata. Given the limited data provided by Netflix—title name, 
    total views, and hours watched—the program seamlessly integrates with the TMDb API to enrich the dataset with genres, production countries, and other essential details. 

    The workflow is fully automated and designed to deliver reliable, well-presented insights. It begins with downloading and storing raw data in JSON format, followed by 
    normalization and a robust analytical process. This culminates in the creation of an organized and visually engaging PowerPoint presentation, offering a clear and 
    comprehensive overview of the findings. 

    Key features of this analysis include:  
    - Identification of the **50 most and least-watched titles**, separated into movies and series.  
    - Detailed visualizations of genre popularity, including **bar charts for total views** and **pie charts for percentage shares**, presented separately for movies and series.  
    - Insights into the popularity of production countries, supported by **heatmaps** and bar charts, categorized by type of content.  
    - Examination of trends in viewing habits relative to content duration, highlighting durations associated with higher or lower viewership.  
    - Comprehensive exploration of the overall popularity of genres and production countries, incorporating both movies and series into unified visualizations.

    By simply providing the Netflix Data Dump for the desired year—starting from the second half of 2023—the program delivers a polished and up-to-date presentation. The results 
    are designed to be not only visually appealing but also insightful for understanding trends and patterns within Netflix's vast library.  

    It is important to note that the program's functionality relies on Netflix maintaining its current format for data sharing. If this format changes in the future, adjustments 
    to the code may be required, particularly in the data selection process. However, as long as the format remains consistent, this tool provides a powerful and fully automated 
    solution for analyzing Netflix's content library efficiently.
    '''

    p1_tp.font.name = 'Arial Narrow'
    p1_tp.font.size = Pt(16)
    p1_tp.font.color.rgb = RGBColor(211, 211, 211)
    p1_tp.alignment = PP_ALIGN.LEFT
    gm_p.text_frame.auto_size = True

    page_p.shapes.add_picture(FOLDER_PRESENTATION + 'TMDB.png' , left= Inches(13.5) , top= Inches(6.5) , width = Inches(2) , height = Inches(2))

    page_p.shapes.add_picture(FOLDER_PRESENTATION + 'NETFLIX.png' , left= Inches(0.5) , top= Inches(6.5) , width = Inches(1) , height = Inches(2))




    # PAGINA 2 NETFLIX
    # PAGE 2 NETFLIX
    page_2 = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    page_2.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_1.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)


    page_2_fulltext = ' Netflix publishes semiannual reports detailing the number of views and hours of playback of its \n audiovisual content. These reports provide basic information on the popularity of movies and TV \n shows, but they only include the title names, the number of views, and the hours played. You can \n check this data at: \n\n Second half of 2023 report: https://about.netflix.com/es/news/what-we-watched-the-second-half-of-2023?utm_source=chatgpt.com \n\n First half of 2024 report: https://about.netflix.com/es/news/what-we-watched-the-first-half-of-2024?utm_source=chatgpt.com \n\n For more recent versions, you can search Google for ''Netflix data dump'' or ''What we watch'' followed \n by the year you''re interested in  \n\n\n Disclaimer: The functionality of the code depends on Netflix continuing to provide the information in \n the same manner as it does in 2024. If the format or structure of the data delivery changes in the \n future, the code may require a rewrite in the section responsible for data selection.'

    gm_2_1 = page_2.shapes.add_textbox(Inches(0.3) , Inches(3) , Inches(5) , Inches(1))
    p2_t1 = gm_2_1.text_frame.paragraphs[0]
    p2_t1.text = page_2_fulltext
    p2_t1.font.name = 'Calibri'
    p2_t1.font.size = Pt(20)
    p2_t1.font.color.rgb = RGBColor(211, 211, 211)
    gm_2_1.text_frame.auto_size = True


    page_2.shapes.add_picture(FOLDER_PRESENTATION + 'NETFLIX.png' , left= Inches(0.5) , top= Inches(0.5) , width = Inches(1) , height = Inches(2))



    # PAGINA 3 TMDb
    # PAGE 3 TMDb
    page_3 = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    page_3.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_1.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)


    page_3_fulltext = 'I would like to express my sincere gratitude to **The Movie Database (TMDb)** and its API\nfor allowing me to use their valuable services in the development of my application.\nThe integration of the TMDb API has been crucial in enhancing the functionality and\nuser experience by providing access to an extensive and up-to-date database of\naudiovisual content.\nThis tool has made it easier to incorporate relevant information about movies, TV shows,\nand other content, significantly enriching the project.\n\nI deeply appreciate the work TMDb does in maintaining such a comprehensive and\naccessible platform, and I acknowledge the positive impact it has had on the performance\nand quality of my application.\nI am very grateful for the opportunity to use their service.\n\nFor more information, you can visit their official website at:\nhttps://www.themoviedb.org/\n\nThis service uses TMDB and the TMDB APIs but is not endorsed or certified by TMDB'

    gm_3_1 = page_3.shapes.add_textbox(Inches(4.7) , Inches(3) , Inches(11) , Inches(1))
    p3_t1 = gm_3_1.text_frame.paragraphs[0]
    p3_t1.text = page_3_fulltext
    p3_t1.font.name = 'Calibri'
    p3_t1.font.size = Pt(20)
    p3_t1.font.color.rgb = RGBColor(211, 211, 211)
    p3_t1.alignment = PP_ALIGN.RIGHT
    gm_3_1.text_frame.auto_size = True

    page_3.shapes.add_picture(FOLDER_PRESENTATION + 'TMDB.png' , left= Inches(13.5) , top= Inches(0.5) , width = Inches(2) , height = Inches(2))



    # PAGINA 4 INDEX
    # PAGE 4 INDEX
    conn = sqlite3.connect(PATH + 'DataBase_normalized.sqlite')
    cur = conn.cursor()

    cur.execute('''
        SELECT
            Information.title
        FROM
            Information
    ''')
    totals = cur.fetchall()
    total = len(totals)
    conn.close()
    total = f'{abs(total):,}'


    page_4 = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    page_4.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_1.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)


    temp_gm = page_4.shapes.add_textbox(Inches(8.7) , Inches(0) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]

    temp_text.text = 'INDEX'

    temp_text.font.name = 'times New Roman'
    temp_text.font.size = Pt(160)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.LEFT
    temp_gm.text_frame.auto_size = True


    temp_gm = page_4.shapes.add_textbox(Inches(12.52) , Inches(0.2) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]

    temp_text.text = 'Trend Analysis of Netflix Content'

    temp_text.font.name = 'times New Roman'
    temp_text.font.size = Pt(16)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.LEFT
    temp_gm.text_frame.auto_size = True



    temp_gm = page_4.shapes.add_textbox(Inches(8.7) , Inches(0.2) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]

    temp_text.text = f'Titles analyzed: {total}'

    temp_text.font.name = 'times New Roman'
    temp_text.font.size = Pt(16)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.LEFT
    temp_gm.text_frame.auto_size = True



    temp_gm = page_4.shapes.add_textbox(Inches(10) , Inches(8.2) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]

    temp_text.text = 'The voting data reflects a public score provided by TMDb,\nnot Netflix. For this reason, it is not considered in the popularity measurements.'

    temp_text.font.name = 'times New Roman'
    temp_text.font.size = Pt(12)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.RIGHT
    temp_gm.text_frame.auto_size = True







    # LISTA PARA RECUPERAR LAS LOS ENLACES A LAS PAGINAS.
    # LIST TO RECOVER LINK TO THE PAGES.
    index_list = list()



    PAGE = 5
    
    # DEFINICION DE LA FUNCION QUE SE ENCARGA DE LA PAGINACION.
    # DEFINITION OF THE FUNCTION RESPONSIBLE FOR PAGINATION.
    def COUNT (NAME_PAGE):
        global PAGE

        gm = NAME_PAGE.shapes.add_textbox(Inches(12.7) , Inches(0) , Inches(6) , Inches(1))
        text = gm.text_frame.paragraphs[0]

        text.text = str(PAGE)
        
        text.font.name = 'Impact'
        text.font.size = Pt(18)
        text.font.color.rgb = RGBColor(211, 211, 211)
        text.alignment = PP_ALIGN.CENTER
        gm.text_frame.auto_size = True



    ############################################################################################################################################################ DEFINITION FOR MOVIES AND SERIES
    ############################################################################################################################################################ DEFINITION FOR MOVIES AND SERIES

    def TV_AND_MOVIES(TYPE , CORRECT_FOLDER):
        global PAGE

        if TYPE == 'tv':
            conn = sqlite3.connect(PATH + 'Analysis_tv.sqlite')
            cur = conn.cursor()

        if TYPE == 'movie':
            conn = sqlite3.connect(PATH + 'Analysis_Movie.sqlite')
            cur = conn.cursor()
        


        ################################################################################################################################################## TOP50s INDENT DEFINITION
        ################################################################################################################################################## TOP50s INDENT DEFINITION
        # TODAS LAS PAGINAS PARA TODOS LOS TOP50
        # ALL PAGES FOR ALL TOP50
        def TOP50s(TYPE , TOP50):
            global PAGE



            # PAGINA  TITULO PARA LOS TOP 50
            # PAGINA  TITLE FOR TOP50s
            PAGE += 1
            temp_page = 'page' + str(PAGE)

            temp_gm1 = 'gm1_page' + str(PAGE)
            temp_gm2 = 'gm2_page' + str(PAGE)

            temp_text1 = 'text1_page' + str(PAGE)
            temp_text2 = 'text2_page' + str(PAGE)



            temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
            temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_3.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
            COUNT(temp_page)



            temp_gm1 = temp_page.shapes.add_textbox(Inches(5) , Inches(2) , Inches(6) , Inches(1))
            temp_text1 = temp_gm1.text_frame.paragraphs[0]


            if TYPE == 'tv': 
                temp_text1.text = 'Analysis of \nTv/Series Content \nTop' + TOP50
            
            if TYPE == 'movie': 
                temp_text1.text = 'Analysis of \nMovies Content \nTop' + TOP50


            temp_text1.font.name = 'Impact'
            temp_text1.font.size = Pt(64)
            temp_text1.font.color.rgb = RGBColor(211, 211, 211)
            temp_text1.alignment = PP_ALIGN.CENTER
            temp_gm1.text_frame.auto_size = True



            temp_gm2 = temp_page.shapes.add_textbox(Inches(5) , Inches(5) , Inches(6) , Inches(1))
            temp_text2 = temp_gm2.text_frame.paragraphs[0]


            if TOP50 == '50_Most_Viewed':
                temp_text2.text = """
                This analysis highlights the 50 most-viewed titles on Netflix, offering insights into the platform's top-performing content.
                These selections reflect audience preferences and showcase the titles that have captivated viewers worldwide.
                """
            if TOP50 == '50_Most_Watched_by_Time_Played':
                temp_text2.text = """
                This analysis highlights the 50 most-watched titles on Netflix by hours played, showcasing the content that keeps audiences engaged the longest.
                These titles reveal trends in viewer commitment and platform engagement.
                """
            if TOP50 == '50_Less_Viewed':
                temp_text2.text = """
                This analysis highlights the 50 least-viewed titles on Netflix, offering a unique perspective on content that struggles to capture audience attention.
                These selections provide insights into underperforming trends on the platform.
                """
            if TOP50 == '50_Less_Watched_by_Time_Played':
                temp_text2.text = """
                This analysis highlights the 50 least-watched titles on Netflix by hours played, focusing on content with the lowest engagement.
                These titles reflect minimal viewer commitment and provide valuable context on platform trends.
                """


            temp_text2.font.name = 'Calibri'
            temp_text2.font.size = Pt(18)
            temp_text2.font.color.rgb = RGBColor(211, 211, 211)
            temp_text2.alignment = PP_ALIGN.CENTER
            temp_gm2.text_frame.auto_size = True


            index_list.append(PAGE)



            # RECUPERECION DE LOS DATOS, PARA CADA CASO.
            # RECOVERY OF THE DATA, FOR EVERY CASE.
            
            if TOP50 == '50_Most_Viewed':
                cur.execute('''
                    SELECT
                        Best50.views,
                        Best50.titles,
                        Best50.runtime,
                        Best50.votes,
                        Best50.hours_viewed,
                        Best50.countries,
                        Best50.genres,
                        Best50.posters
                    FROM
                        Best50
                ''')

            if TOP50 == '50_Most_Watched_by_Time_Played':
                cur.execute('''
                    SELECT
                        Best50_Hours.views,
                        Best50_Hours.titles,
                        Best50_Hours.runtime,
                        Best50_Hours.votes,
                        Best50_Hours.hours_viewed,
                        Best50_Hours.countries,
                        Best50_Hours.genres,
                        Best50_Hours.posters
                    FROM
                        Best50_Hours
                ''')
            
            if TOP50 == '50_Less_Viewed':
                cur.execute('''
                    SELECT
                        Worst50.views,
                        Worst50.titles,
                        Worst50.runtime,
                        Worst50.votes,
                        Worst50.hours_viewed,
                        Worst50.countries,
                        Worst50.genres,
                        Worst50.posters
                    FROM
                        Worst50
                ''')
            if TOP50 == '50_Less_Watched_by_Time_Played':
                cur.execute('''
                    SELECT
                        Worst50_Hours.views,
                        Worst50_Hours.titles,
                        Worst50_Hours.runtime,
                        Worst50_Hours.votes,
                        Worst50_Hours.hours_viewed,
                        Worst50_Hours.countries,
                        Worst50_Hours.genres,
                        Worst50_Hours.posters
                    FROM
                        Worst50_Hours
                ''')


            Tops50_data = cur.fetchall()


            # DESGLOSE DE LOS DATOS.
            # BREAKDOWN OF THE DATA.
            views_temp = list()
            title_ = list()
            runtime_ = list()
            votes_ = list()
            hours_temp = list()
            countries_ = list()
            genres_ = list()
            posters_ = list() 

            for row in Tops50_data:
                views_temp.append(row[0])
                title_.append(row[1])
                runtime_.append(row[2])
                votes_.append(row[3])
                hours_temp.append(row[4])
                countries_.append(row[5])
                genres_.append(row[6])
                posters_.append(row[7])
            
            views_ = list()
            hours_ = list()
            for v , h in zip(views_temp , hours_temp):
                h = int(h)
                vi = f'{abs(v):,}'
                ho = f'{abs(h):,}'
                views_.append(vi)
                hours_.append(ho)
            


            ######################## CREACION DE LA CARPETA QUE CONTENDRA LOS POSTERS 
            ######################## CREATION OF THE FOLDER THAT WILL CONTAIN THE POSTERS.
            if not os.path.exists (PATH + TYPE + '_folder'):
                os.mkdir(PATH + TYPE + '_folder')

            FOLDER_POSTERS = (PATH + TYPE + '_folder' + '//')


            # DESCARGA DE LOS PORTES
            # POSTERS DOWNLOADED
            image_name_counter = 0
            for row in posters_:
                image_name_counter +=1
                
                if os.path.exists(FOLDER_POSTERS + TOP50 + '_' + str(image_name_counter) + '.png'):
                    continue

                try:
                    image_url = 'https://image.tmdb.org/t/p/w500' + row
                    file_name = FOLDER_POSTERS + TOP50 + '_' + str(image_name_counter) + '.png'

                    image_requested = requests.get(image_url)
                except:
                    print('Was Not Posible to Download One Poster Image')
                    continue

                if image_requested.status_code ==200:
                    with open(file_name , 'wb') as file:
                        file.write(image_requested.content)
                else:
                    print('Was Not Posible to Download One Poster Image')




            ########## BUCLE PARA ACOMODAR CADA TITULO
            ########## LOOP FOR PLACE EVERY TITLE

            next_page = 25
            image_name_counter = 0
            for vi , ti , ru , vo , ho , co , ge , po in zip(views_ , title_ , runtime_ , votes_ , hours_ , countries_ , genres_ , range(50)):

                # FILTRO DE TITULOS DEMACIADO LARGOS.
                # FILTER OF TITLES THAT ARE TOO LONG.
                if len(ti) >= 45:
                    ti = ti[ :46]



                # CREACION DE UNA NUVA PAGINA PARA LA SEGUNDA MITAD
                # CREATION OF A NEW PAGE FOR DE SECOND HALF
                if next_page >= 25:
                    next_page = 0
                    PAGE += 1
                    temp_page = 'page' + str(PAGE)

                    TEXT = 0

                    temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
                    temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_3.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
                    COUNT(temp_page)


                    # ASIGNACION DE VARIABLES A LA UBICACION DEL TEXTO Y LAS IMAGENES.
                    # ASIGNATION OF VARIABLES FOR TEXT AND IMAGE PLACE.
                    img_left = Inches(0.6)
                    img_top = Inches(1)
                    img_width = Inches(0.7)
                    img_height = Inches(1)

                    txt_left = img_left + img_width + Inches(0.1)
                    txt_top = img_top - Inches(0.13)
                    txt_width = Inches(1)
                    txt_height = img_height


                    jump_line = 0
                


                # CONTADOR PARA CREAR LA SEGUNDA PAGINA. Y CONTADOR SELECTOR DE IMAGENES.
                # COUNTER TO CREATE DE SECOND PAGE. AND SELECT IMAGE COUNTER.
                next_page += 1  
                image_name_counter += 1



                try:
                    # IMAGEN
                    # IMAGE
                    image_name = FOLDER_POSTERS + TOP50 + '_' + str(image_name_counter) + '.png'
                    temp_page.shapes.add_picture( image_name , left= img_left , top= img_top , width = img_width , height = img_height)
                except:
                    # IMAGEN NEGRA
                    # BLACK IMAGE
                    image_name = FOLDER_PRESENTATION + str(0) + '.png'
                    temp_page.shapes.add_picture( image_name , left= img_left , top= img_top , width = img_width , height = img_height)
                


                # INFORMACION DE LA PELICULA.
                # MOVIE INFORMATION
                TEXT += 1
                loop_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
                loop_text = 'text' + str(TEXT) + '_page' + str(PAGE)
                

                loop_gm = temp_page.shapes.add_textbox(txt_left , txt_top , txt_width , txt_height)
                loop_text = loop_gm.text_frame.paragraphs[0]

                loop_text.text = '\nViews: ' + str(vi) + '\nTitle: ' + str(ti) + '\nRuntime: ' + str(ru) + '\nVotes: ' + str(vo) + '\nReproducted Hours :' + str(ho) + '\nCountrie: ' + str(co) + '\nGenres: ' + str(ge)

                loop_text.font.name = 'Calibri'
                loop_text.font.size = Pt(8)
                loop_text.font.color.rgb = RGBColor(211, 211, 211)
                loop_text.alignment = PP_ALIGN.LEFT
                loop_gm.text_frame.auto_size = True


                # NUMERO DE LUGAR.
                # PLACE NUMBER.
                TEXT += 1
                loop_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
                loop_text = 'text' + str(TEXT) + '_page' + str(PAGE)
                

                loop_gm = temp_page.shapes.add_textbox(img_left , img_top , txt_width , txt_height)
                loop_text = loop_gm.text_frame.paragraphs[0]

                loop_text.text = str(po + 1)

                loop_text.font.name = 'Impact'
                loop_text.font.size = Pt(16)
                loop_text.font.color.rgb = RGBColor(255, 255, 255)
                loop_text.alignment = PP_ALIGN.LEFT
                loop_gm.text_frame.auto_size = True



                # REORGANIZACION DE LO DATOS DENTRO DE LA DIAPOSITIVA. HORISONTAL
                # REORGANIZATION OF THE DATA WITHIN THE SLIDE. HORIZONTAL
                img_left = img_left + img_width + Inches(0.1) + Inches(2) + Inches(0.1)
                txt_left = img_left + img_width + Inches(0.1)


                jump_line += 1
                if jump_line >= 5:
                    # REORGANIZACION DE LO DATOS DENTRO DE LA DIAPOSITIVA. VERTICAL
                    # REORGANIZATION OF THE DATA WITHIN THE SLIDE. VERTICAL
                    img_left = Inches(0.6)
                    txt_left = img_left + img_width + Inches(0.1)

                    img_top = img_top + img_height + Inches(0.5)
                    txt_top = img_top - Inches(0.13)

                    jump_line = 0

            

        ################################################################################################################################################## TOP50s INDENT DEFINITION END
        ################################################################################################################################################## TOP50s INDENT DEFINITION END
                
        # USO DE LA FUNCION INDENTADA PARA TODOS LOS TOP50.
        # USE OF THE INDENTED FUNCTION FOR ALL TOP50.

        TOP50s( TYPE , '50_Most_Viewed')

        TOP50s( TYPE , '50_Most_Watched_by_Time_Played')

        TOP50s( TYPE , '50_Less_Viewed')

        TOP50s( TYPE , '50_Less_Watched_by_Time_Played')



        ########## DIAPOSITIVAS DESTINADAS A LA PRESENTACION DE LOS GRAFICOS DE LOS GENEROS.
        ########## SLIDES INTENDED FOR THE PRESENTATION OF THE GENRE GRAPHS.
        PAGE += 1
        temp_page = 'page' + str(PAGE)

        TEXT = 0
        temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
        temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


        temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
        temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_2.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
        COUNT(temp_page)


        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(5) , Inches(2) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]


        if TYPE == 'tv': 
            temp_text.text = 'Analysis of \nTv/Series Content \nGenres Popularity'
        
        if TYPE == 'movie': 
            temp_text.text = 'Analysis of \nMovies Content \nGenres Popularity'


        temp_text.font.name = 'Impact'
        temp_text.font.size = Pt(64)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.CENTER
        temp_gm.text_frame.auto_size = True



        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(5) , Inches(5) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]

        temp_text.text = '''
    This analysis explores the popularity of genres across Netflix's entire catalog, using comprehensive viewership data.
    A bar chart highlights the total number of views for each genre, providing a clear comparison of audience preferences.
    Additionally, a pie chart presents the percentage share of views, offering a visual breakdown of each genre's popularity relative to the platform as a whole.
        '''
        temp_text.font.name = 'Calibri'
        temp_text.font.size = Pt(18)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.CENTER
        temp_gm.text_frame.auto_size = True



        index_list.append(PAGE)



        # OPTENCION DEL DATO DE MARGEN DE ERROR PARA LOS GENEROS.
        # OBTAINING THE ERROR MARGIN DATA FOR THE GENRES
        cur.execute('''
            SELECT
                Margin_Of_Error.null_genres,
                Margin_Of_Error.margin_genres
            FROM
                Margin_Of_Error
        ''')
        genres_errors = cur.fetchall()

        null_titles_genres = str(genres_errors[0][0])
        error_margin_genres = str(genres_errors[0][1])
        error_margin_genres = error_margin_genres[ :5] + ' %'

        

        # GRAFICO DE BARRAS PARA LOS GENEROS.
        # BAR CHART FOR THE GENRES.
        PAGE += 1
        temp_page = 'page' + str(PAGE)

        TEXT = 0
        temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
        temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


        temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
        temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_2.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
        COUNT(temp_page)


        if TYPE == 'tv':
            temp_page.shapes.add_picture(CORRECT_FOLDER + 'Tv_Graphs_genres_bar.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
        
        if TYPE == 'movie':
            temp_page.shapes.add_picture(CORRECT_FOLDER + 'Movie_Graphs_genres_bar.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)



        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(0) , Inches(0) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]

        temp_text.text = f'''
        This bar chart shows the popularity of genres based on their total views. 
        There is a small error margin due to titles with missing genre data.
        (This margin varies depending on the analyzed Data Dump.)
        Titles With No Genres: {null_titles_genres}
        Error Margin: {error_margin_genres}
        '''

        temp_text.font.name = 'times New Roman'
        temp_text.font.size = Pt(14)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.LEFT
        temp_gm.text_frame.auto_size = True
        


        # GRAFICOS DE TORTA PARA LOS GENEROS.
        # PIE CHARTS FOR THE GENRES.
        PAGE += 1
        temp_page = 'page' + str(PAGE)

        TEXT = 0
        temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
        temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


        temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
        temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_2.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
        COUNT(temp_page)


        if TYPE == 'tv':
            images_ = list()
            images_.append('Tv_Graphs_genres_full_pie.png')
            images_.append('Tv_Graphs_genres_pie_h1.png')
            images_.append('Tv_Graphs_genres_pie_h2.png')
            correct_width = dict ()
            correct_height = dict ()

            helper_counter = 0
            for img in images_:
                helper_counter += 1
                with Image.open(CORRECT_FOLDER + img) as IMG:
                    pixel_width , pixel_height = IMG.size

                inches_original_w = pixel_width / 96
                inches_original_h = pixel_height / 96

                if helper_counter == 1:
                    max_w = 9.4
                    max_h = 6
                else:
                    max_w = 6
                    max_h = 4


                ratio = min(max_w / inches_original_w , max_h / inches_original_h)

                correct_width[helper_counter] = inches_original_w * ratio
                correct_height[helper_counter] = inches_original_h * ratio



            temp_page.shapes.add_picture(CORRECT_FOLDER + 'Tv_Graphs_genres_full_pie.png' , left= Inches(0.2) , top= Inches(3) , width=Inches(correct_width[1]) , height=Inches(correct_height[1]))

            temp_page.shapes.add_picture(CORRECT_FOLDER + 'Tv_Graphs_genres_pie_h1.png' , left= Inches(9.8) , top= Inches(0.5) , width=Inches(correct_width[2]) , height=Inches(correct_height[2]))
            temp_page.shapes.add_picture(CORRECT_FOLDER + 'Tv_Graphs_genres_pie_h2.png' , left= Inches(9.8) , top= Inches(4.8) , width=Inches(correct_width[3]) , height=Inches(correct_height[3]))
        


        if TYPE == 'movie':
            images_ = list()
            images_.append('Movie_Graphs_genres_full_pie.png')
            images_.append('Movie_Graphs_genres_pie_h1.png')
            images_.append('Movie_Graphs_genres_pie_h2.png')
            correct_width = dict ()
            correct_height = dict ()

            helper_counter = 0
            for img in (images_):
                helper_counter += 1
                with Image.open(CORRECT_FOLDER + img) as IMG:
                    pixel_width , pixel_height = IMG.size

                inches_original_w = pixel_width / 96
                inches_original_h = pixel_height / 96

                if helper_counter == 1:
                    max_w = 9.4
                    max_h = 6
                else:
                    max_w = 6
                    max_h = 4

                ratio = min(max_w / inches_original_w , max_h / inches_original_h)

                correct_width[helper_counter] = inches_original_w * ratio
                correct_height[helper_counter] = inches_original_h * ratio



            temp_page.shapes.add_picture(CORRECT_FOLDER + 'Movie_Graphs_genres_full_pie.png' , left= Inches(0.2) , top= Inches(3) , width=Inches(correct_width[1]) , height=Inches(correct_height[1])) 

            temp_page.shapes.add_picture(CORRECT_FOLDER + 'Movie_Graphs_genres_pie_h1.png' , left= Inches(9.8) , top= Inches(0.5) , width=Inches(correct_width[2]) , height=Inches(correct_height[2])) 
            temp_page.shapes.add_picture(CORRECT_FOLDER + 'Movie_Graphs_genres_pie_h2.png' , left= Inches(9.8) , top= Inches(4.8) , width=Inches(correct_width[3]) , height=Inches(correct_height[3])) 



        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(0) , Inches(0) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]

        temp_text.text = f'''
        This pie charts shows the popularity of genres based on their total views. 
        There is a full one, and to half of the same information cut in two. 
        There is a small error margin due to titles with missing genre data.
        (This margin varies depending on the analyzed Data Dump.)
        Titles With No Genres: {null_titles_genres}
        Error Margin: {error_margin_genres}
        '''

        temp_text.font.name = 'times New Roman'
        temp_text.font.size = Pt(14)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.LEFT
        temp_gm.text_frame.auto_size = True



        ########## COMIENZO DE LAS DIAPOSITIVAS QUE PRESENTARÁN LOS PAÍSES PRODUCTORES.
        ########## BEGINNING OF THE SLIDES PRESENTING THE PRODUCING COUNTRIES.


        # OPTENCION DEL DATO DE MARGEN DE ERROR PARA LOS PAISES.
        # OBTAINING THE ERROR MARGIN DATA FOR COUNTRIES
        cur.execute('''
            SELECT
                Margin_Of_Error.null_countries,
                Margin_Of_Error.margin_countries
            FROM
                Margin_Of_Error
        ''')
        countries_errors = cur.fetchall()

        null_titles_countries = str(countries_errors[0][0])
        error_margin_countries = str(countries_errors[0][1])
        error_margin_countries = error_margin_countries[ :5] + ' %'


        # PAGINA DEL TITULO
        # TITLE PAGE
        PAGE += 1
        temp_page = 'page' + str(PAGE)

        TEXT = 0
        temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
        temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


        temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
        temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_2.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
        COUNT(temp_page)


        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(5) , Inches(2) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]


        if TYPE == 'tv': 
            temp_text.text = 'Analysis of \nTv/Series Content \nProducing Countries Popularity'
        
        if TYPE == 'movie': 
            temp_text.text = 'Analysis of \nMovies Content \nProducing Countries Popularity'


        temp_text.font.name = 'Impact'
        temp_text.font.size = Pt(64)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.CENTER
        temp_gm.text_frame.auto_size = True



        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(5) , Inches(5) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]

        temp_text.text = '''
    This analysis examines the popularity of producer countries in Netflix content, highlighting how productions from different 
    countries impact the platform's global audience. Viewing trends and the impact of producer countries on user preferences 
    will be analyzed.
        '''
        temp_text.font.name = 'Calibri'
        temp_text.font.size = Pt(18)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.CENTER
        temp_gm.text_frame.auto_size = True



        index_list.append(PAGE)

        
        # CONTEO DE LA CANTIDAD DE GRAFICOS GENERADOS PARA LOS PAISES.
        # COUNT OF THE NUMBER OF GRAPHS GENERATED FOR THE COUNTRIES.
        graphs_cuantiti_count = 1


        if TYPE == 'tv':

            while True:

                if (os.path.exists(CORRECT_FOLDER + 'Tv_Graphs_countries_bars' + str(graphs_cuantiti_count) + '.png')):
                    graphs_cuantiti_count += 1

                else:
                    break

        

        if TYPE == 'movie':

            while True:

                if (os.path.exists(CORRECT_FOLDER + 'Movie_Graphs_countries_bars' + str(graphs_cuantiti_count) + '.png')):
                    graphs_cuantiti_count += 1

                else:
                    break



        # BUCLE QUE CREA TODAS LAS DIAPOSITIVAS DEL ANALISIS DE LOS PAISES SEGUN LA CANTIDAD QUE SEA NECESARIA.
        # LOOP THAT CREATES ALL THE SLIDES FOR THE COUNTRY ANALYSIS ACCORDING TO THE REQUIRED QUANTITY
        for num in range(graphs_cuantiti_count - 1):
            num = num + 1

            PAGE += 1
            temp_page = 'page' + str(PAGE)

            TEXT = 0
            temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
            temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


            temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
            temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_4.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
            COUNT(temp_page)



            TEXT += 1
            temp_gm = temp_page.shapes.add_textbox(Inches(0) , Inches(0) , Inches(6) , Inches(1))
            temp_text = temp_gm.text_frame.paragraphs[0]

            temp_text.text = f'''
        This bar and pie charts shows the popularity of the Producing Countries based on their total views.
        Keep in mind that the largest chart on one page is smaller than the last one on the previous page.
        There is a small error margin due to titles with missing countries data.
        (This margin varies depending on the analyzed Data Dump.)
        Titles With No Genres: {null_titles_countries}
        Error Margin: {error_margin_countries}
            '''

            temp_text.font.name = 'times New Roman'
            temp_text.font.size = Pt(14)
            temp_text.font.color.rgb = RGBColor(211, 211, 211)
            temp_text.alignment = PP_ALIGN.LEFT
            temp_gm.text_frame.auto_size = True


            # SE MIDEN EL TAMAÑO DE LA IMAGEN REAL EN PIXELESM LUEGO SE HACE UNA FORMULA PARA PASAR ESE TAMAÑO A POLGADAS Y SE APLICAN LOS DIAGRAMAS DE PIE CON LA PROPORCION CORRECTA.
            # THE ACTUAL IMAGE SIZE IS MEASURED IN PIXELS, THEN A FORMULA IS USED TO CONVERT THAT SIZE TO INCHES, AND PIE CHARTS ARE APPLIED WITH THE CORRECT PROPORTION.
            if TYPE == 'tv':
                with Image.open(CORRECT_FOLDER + 'Tv_Graphs_countries_pies' + str(num) + '.png') as IMG:
                    pixel_width , pixel_height = IMG.size

                inches_original_w = pixel_width / 96
                inches_original_h = pixel_height / 96

                max_w = 6
                max_h = 3

                ratio = min(max_w / inches_original_w , max_h / inches_original_h)

                correct_width = inches_original_w * ratio
                correct_height = inches_original_h * ratio

        

                temp_page.shapes.add_picture(CORRECT_FOLDER + 'Tv_Graphs_countries_bars' + str(num) + '.png' , left=Inches(0) , top=Inches(3) , width=diapo_.slide_width , height=(diapo_.slide_height - (diapo_.slide_height/3) ))

                temp_page.shapes.add_picture(CORRECT_FOLDER + 'Tv_Graphs_countries_pies' + str(num) + '.png' , left= Inches(9.8) , top= Inches(0) , width=Inches(correct_width) , height=Inches(correct_height))
            


            if TYPE == 'movie':
                with Image.open(CORRECT_FOLDER + 'Movie_Graphs_countries_pies' + str(num) + '.png') as IMG:
                    pixel_width , pixel_height = IMG.size

                inches_original_w = pixel_width / 96
                inches_original_h = pixel_height / 96

                max_w = 6
                max_h = 4

                ratio = min(max_w / inches_original_w , max_h / inches_original_h)

                correct_width = inches_original_w * ratio
                correct_height = inches_original_h * ratio



                temp_page.shapes.add_picture(CORRECT_FOLDER + 'Movie_Graphs_countries_bars' + str(num) + '.png' , left=Inches(0) , top=Inches(3) , width=diapo_.slide_width , height=(diapo_.slide_height - (diapo_.slide_height/3) )) 

                temp_page.shapes.add_picture(CORRECT_FOLDER + 'Movie_Graphs_countries_pies' + str(num) + '.png' , left= Inches(9.8) , top= Inches(0) , width=Inches(correct_width) , height=Inches(correct_height)) 



            



        ########## COMIENZO DE LAS DIAPOSITIVAS QUE PRESENTARÁN LA POPULARIDAD DE LOS TIEMPO DE REPRODUCCION..
        ########## BEGINNING OF THE SLIDES PRESENTING THE RUNTIMES POPULARITY.

        # OPTENCION DE LOS DATOS DE LA POPULARIDAD DE LAS VISTAS.
        # OBTAINING VIEW POPULARITY DATA.
        cur.execute('''
            SELECT
                Large_Runtime_Popularity.runtime_average,
                Large_Runtime_Popularity.views_by_average
            FROM
                Large_Runtime_Popularity   
    ''')
        large_runtime_popularity = cur.fetchall()

        runtime_popularuty_list = list()
        for row in large_runtime_popularity:
            views = f'{abs(row[1]):,}'
            runtime_popularuty_list.append(row[0] + ' - ' + views)



        # PAGINA DEL TITULO
        # TITLE PAGE
        PAGE += 1
        temp_page = 'page' + str(PAGE)

        TEXT = 0
        temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
        temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


        temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
        temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_5.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
        COUNT(temp_page)


        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(5) , Inches(2) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]


        if TYPE == 'tv': 
            temp_text.text = 'Analysis of \nTv/Series Content \nRelationship Between Duration \nand Playback Popularity'
        
        if TYPE == 'movie': 
            temp_text.text = 'Analysis of \nMovies Content \nRelationship Between Duration \nand Playback Popularity'


        temp_text.font.name = 'Impact'
        temp_text.font.size = Pt(64)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.CENTER
        temp_gm.text_frame.auto_size = True



        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(5) , Inches(6) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]

        temp_text.text = '''
    This analysis examines the relationship between content duration and its popularity on Netflix. Averages of playback times 
    will be calculated, and the number of views within each interval will be evaluated. The results will be presented in bar and 
    scatter plots to illustrate these trends.
        '''
        temp_text.font.name = 'Calibri'
        temp_text.font.size = Pt(18)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.CENTER
        temp_gm.text_frame.auto_size = True



        index_list.append(PAGE)


        
        # GRAFICO DE LINEAS PARA LOS LOS TIEMPOS DE REPRODUCCION.
        # LINE CHART FOR THE RUNTIMES.
        PAGE += 1
        temp_page = 'page' + str(PAGE)

        TEXT = 0
        temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
        temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


        temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
        temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_5.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
        COUNT(temp_page)


        if TYPE == 'tv':
            temp_page.shapes.add_picture(CORRECT_FOLDER + 'Tv_Graphs_popularity_lines.png' , 0 , top=Inches(1.5) , width=diapo_.slide_width , height=(diapo_.slide_height - (diapo_.slide_height / 6)))
        
        if TYPE == 'movie':
            temp_page.shapes.add_picture(CORRECT_FOLDER + 'Movie_Graphs_popularity_lines.png' , 0 , top=Inches(1.5) , width=diapo_.slide_width , height=(diapo_.slide_height - (diapo_.slide_height / 6)))



        # GRAFICO DE DISPERCION PARA LOS LOS TIEMPOS DE REPRODUCCION.
        # SCATTER CHART FOR THE RUNTIMES.
        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(0) , Inches(0) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]

        temp_text.text = f'''
    This line chart shows how popularity varies according to content 
    duration, identifying trends as playback times change. It 
    provides a clear view of the relationship between duration 
    and the number of views over time.
        '''

        temp_text.font.name = 'times New Roman'
        temp_text.font.size = Pt(14)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.LEFT
        temp_gm.text_frame.auto_size = True


        # GRAFICO DE LINEAS PARA LOS LOS TIEMPOS DE REPRODUCCION.
        # LINE CHART FOR THE RUNTIMES.
        PAGE += 1
        temp_page = 'page' + str(PAGE)

        TEXT = 0
        temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
        temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


        temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
        temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_5.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
        COUNT(temp_page)


        if TYPE == 'tv':
            temp_page.shapes.add_picture(CORRECT_FOLDER + 'Tv_Graphs_popularity_scatter.png' , 0 , top=Inches(1.5) , width=diapo_.slide_width , height=(diapo_.slide_height - (diapo_.slide_height / 6)))
        
        if TYPE == 'movie':
            temp_page.shapes.add_picture(CORRECT_FOLDER + 'Movie_Graphs_popularity_scatter.png' , 0 , top=Inches(1.5) , width=diapo_.slide_width , height=(diapo_.slide_height - (diapo_.slide_height / 6)))



        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(0) , Inches(0) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]

        temp_text.text = f'''
    This scatter plot expands the analysis by showing the relationship between 
    popularity and content duration, but with more reference points. By 
    including more averages, it provides a more detailed and precise 
    view of viewing trends across different time intervals.
        '''

        temp_text.font.name = 'times New Roman'
        temp_text.font.size = Pt(14)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.LEFT
        temp_gm.text_frame.auto_size = True



        # LEYENDA DEL GRAFICO DE DISPERCION.
        # SCATTER CHART LEGEND.
        PAGE += 1
        temp_page = 'page' + str(PAGE)

        TEXT = 0
        temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
        temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


        temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
        temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_5.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
        COUNT(temp_page)



        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(0) , Inches(0) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]

        temp_text.text = f'''
    The legend below details the different duration averages analyzed in the scatter plot, 
    associating each point with a specific time interval to facilitate the interpretation 
    of trends. The number of views for each average corresponds to the accumulated views 
    between the previous time in the list and the represented average. - - Time is shownby Hours/Minutes - -
        '''

        temp_text.font.name = 'times New Roman'
        temp_text.font.size = Pt(14)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.LEFT
        temp_gm.text_frame.auto_size = True



        left_ = 0
        top_ = 1.5
        helper_counter = 0
        
        last_runtime = runtime_popularuty_list[-1]

        for row in runtime_popularuty_list:
            helper_counter += 1

            top_ += 0.5

            TEXT += 1
            temp_gm = temp_page.shapes.add_textbox(Inches(left_) , Inches(top_) , Inches(6) , Inches(1))
            temp_text = temp_gm.text_frame.paragraphs[0]

            temp_text.text = row

            temp_text.font.name = 'Calibri'
            temp_text.font.size = Pt(14)
            temp_text.font.color.rgb = RGBColor(211, 211, 211)
            temp_text.alignment = PP_ALIGN.LEFT
            temp_gm.text_frame.auto_size = True

            if helper_counter >= 13:
                left_ += 2
                top_ = 1.5
                helper_counter = 0
        


        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(10) , Inches(0) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]

        temp_text.text = f'''
    The last average in the list was not included in the charts because it represents 
    extremely long durations with very few views. Its inclusion would distort the 
    tables and make them difficult to read clearly
    -> {last_runtime}
        '''

        temp_text.font.name = 'times New Roman'
        temp_text.font.size = Pt(14)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.RIGHT
        temp_gm.text_frame.auto_size = True
        

        
        conn.close()



    ####################################################################################################################################################### DEFINITION FOR MOVIES AND SERIES END.
    ####################################################################################################################################################### DEFINITION FOR MOVIES AND SERIES END.

    # EJECUCION DE LA DEFINICION TV AND MOVIES.
    # RUNING FOR TV AND MOVIES DEFINITION.
    TV_AND_MOVIES('tv' , FOLDER_TV)

    TV_AND_MOVIES('movie' , FOLDER_MOVIE)






    ############### COMIENZO DE LAS DIAPOSITIVAS DEL ANALICIS GLOVAL, EL CUAL ANALIZA TANTO PELICULAS COMO RESRIES.
    ############### BEGINNING OF THE SLIDES FOR THE GLOBAL ANALYSIS, WHICH COVERS BOTH MOVIES AND SERIES.


    # CONECCION CON EL ANALISIS.
    # ANALYSIS CONECTION.
    conn = sqlite3.connect(PATH + 'Analysis_Global.sqlite')
    cur = conn.cursor()



    # DIAPOSITIVA DE INDUCCION AL ANALISIS GLOBAL.
    # INTRODUCTORY SLIDE TO THE GLOBAL ANALYSIS.
    PAGE += 1
    temp_page = 'page' + str(PAGE)

    TEXT = 0
    temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
    temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


    temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_6.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
    COUNT(temp_page)


    TEXT += 1
    temp_gm = temp_page.shapes.add_textbox(Inches(5) , Inches(2) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]



    temp_text.text = 'Analysis of \nGlobal Content Popularity'



    temp_text.font.name = 'Impact'
    temp_text.font.size = Pt(64)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.CENTER
    temp_gm.text_frame.auto_size = True



    TEXT += 1
    temp_gm = temp_page.shapes.add_textbox(Inches(5) , Inches(5) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]

    temp_text.text = '''
    From this point onward, the information presented and the analyses conducted 
    will be based on data from both movies and series, aiming to conclude with a global 
    analysis that offers us broader insight. This analysis will allow us to explore 
    the overall behavior of the content on the platform in terms of popularity.
    '''
    temp_text.font.name = 'Calibri'
    temp_text.font.size = Pt(24)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.CENTER
    temp_gm.text_frame.auto_size = True



    index_list.append(PAGE)



    ########## DIAPOSITIVAS DESTINADAS A LA PRESENTACION DE LOS GRAFICOS DE LOS GENEROS.
    ########## SLIDES INTENDED FOR THE PRESENTATION OF THE GENRE GRAPHS.
    PAGE += 1
    temp_page = 'page' + str(PAGE)

    TEXT = 0
    temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
    temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


    temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_6.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
    COUNT(temp_page)


    TEXT += 1
    temp_gm = temp_page.shapes.add_textbox(Inches(5) , Inches(2) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]



    temp_text.text = 'Analysis of \nGlobal Content \nGenres Popularity'




    temp_text.font.name = 'Impact'
    temp_text.font.size = Pt(64)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.CENTER
    temp_gm.text_frame.auto_size = True



    TEXT += 1
    temp_gm = temp_page.shapes.add_textbox(Inches(5) , Inches(5) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]

    temp_text.text = '''
    This analysis explores the popularity of genres across Netflix's entire catalog, using comprehensive viewership data.
    A bar chart highlights the total number of views for each genre, providing a clear comparison of audience preferences.
    Additionally, a pie chart presents the percentage share of views, offering a visual breakdown of each genre's popularity relative to the platform as a whole.
    '''
    temp_text.font.name = 'Calibri'
    temp_text.font.size = Pt(18)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.CENTER
    temp_gm.text_frame.auto_size = True



    index_list.append(PAGE)



    # OPTENCION DEL DATO DE MARGEN DE ERROR PARA LOS GENEROS.
    # OBTAINING THE ERROR MARGIN DATA FOR THE GENRES
    cur.execute('''
        SELECT
            Margin_Of_Error.null_genres,
            Margin_Of_Error.margin_genres
        FROM
            Margin_Of_Error
    ''')
    genres_errors = cur.fetchall()

    null_titles_genres = str(genres_errors[0][0])
    error_margin_genres = str(genres_errors[0][1])
    error_margin_genres = error_margin_genres[ :5] + ' %'



    # GRAFICO DE BARRAS PARA LOS GENEROS.
    # BAR CHART FOR THE GENRES.
    PAGE += 1
    temp_page = 'page' + str(PAGE)

    TEXT = 0
    temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
    temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


    temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_6.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
    COUNT(temp_page)


    temp_page.shapes.add_picture(FOLDER_GLOBAL + 'Global_Graphs_genres_bar.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)



    TEXT += 1
    temp_gm = temp_page.shapes.add_textbox(Inches(0) , Inches(0) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]

    temp_text.text = f'''
    This bar chart shows the popularity of genres based on their total views. 
    There is a small error margin due to titles with missing genre data.
    (This margin varies depending on the analyzed Data Dump.)
    Titles With No Genres: {null_titles_genres}
    Error Margin: {error_margin_genres}
    '''

    temp_text.font.name = 'times New Roman'
    temp_text.font.size = Pt(14)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.LEFT
    temp_gm.text_frame.auto_size = True



    # GRAFICOS DE TORTA PARA LOS GENEROS.
    # PIE CHARTS FOR THE GENRES.
    PAGE += 1
    temp_page = 'page' + str(PAGE)

    TEXT = 0
    temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
    temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


    temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_6.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
    COUNT(temp_page)



    images_ = list()
    images_.append('Global_Graphs_genres_full_pie.png')
    images_.append('Global_Graphs_genres_pie_h1.png')
    images_.append('Global_Graphs_genres_pie_h2.png')
    correct_width = dict ()
    correct_height = dict ()

    helper_counter = 0
    for img in images_:
        helper_counter += 1
        with Image.open(FOLDER_GLOBAL + img) as IMG:
            pixel_width , pixel_height = IMG.size

        inches_original_w = pixel_width / 96
        inches_original_h = pixel_height / 96

        if helper_counter == 1:
            max_w = 9.4
            max_h = 6
        else:
            max_w = 6
            max_h = 4


        ratio = min(max_w / inches_original_w , max_h / inches_original_h)

        correct_width[helper_counter] = inches_original_w * ratio
        correct_height[helper_counter] = inches_original_h * ratio



    temp_page.shapes.add_picture(FOLDER_GLOBAL + 'Global_Graphs_genres_full_pie.png' , left= Inches(0.2) , top= Inches(3) , width=Inches(correct_width[1]) , height=Inches(correct_height[1]))

    temp_page.shapes.add_picture(FOLDER_GLOBAL + 'Global_Graphs_genres_pie_h1.png' , left= Inches(9.8) , top= Inches(0.5) , width=Inches(correct_width[2]) , height=Inches(correct_height[2]))
    temp_page.shapes.add_picture(FOLDER_GLOBAL + 'Global_Graphs_genres_pie_h2.png' , left= Inches(9.8) , top= Inches(4.8) , width=Inches(correct_width[3]) , height=Inches(correct_height[3]))



    TEXT += 1
    temp_gm = temp_page.shapes.add_textbox(Inches(0) , Inches(0) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]

    temp_text.text = f'''
    This pie charts shows the popularity of genres based on their total views. 
    There is a full one, and to half of the same information cut in two. 
    There is a small error margin due to titles with missing genre data.
    (This margin varies depending on the analyzed Data Dump.)
    Titles With No Genres: {null_titles_genres}
    Error Margin: {error_margin_genres}
    '''

    temp_text.font.name = 'times New Roman'
    temp_text.font.size = Pt(14)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.LEFT
    temp_gm.text_frame.auto_size = True



    ########## COMIENZO DE LAS DIAPOSITIVAS QUE PRESENTARÁN LOS PAÍSES PRODUCTORES.
    ########## BEGINNING OF THE SLIDES PRESENTING THE PRODUCING COUNTRIES.


    # OPTENCION DEL DATO DE MARGEN DE ERROR PARA LOS PAISES.
    # OBTAINING THE ERROR MARGIN DATA FOR COUNTRIES
    cur.execute('''
        SELECT
            Margin_Of_Error.null_countries,
            Margin_Of_Error.margin_countries
        FROM
            Margin_Of_Error
    ''')
    countries_errors = cur.fetchall()

    null_titles_countries = str(countries_errors[0][0])
    error_margin_countries = str(countries_errors[0][1])
    error_margin_countries = error_margin_countries[ :5] + ' %'


    # PAGINA DEL TITULO
    # TITLE PAGE
    PAGE += 1
    temp_page = 'page' + str(PAGE)

    TEXT = 0
    temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
    temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


    temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_7.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
    COUNT(temp_page)


    TEXT += 1
    temp_gm = temp_page.shapes.add_textbox(Inches(5) , Inches(2) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]


    temp_text.text = 'Analysis of \nGlobal Content \nProducing Countries Popularity'


    temp_text.font.name = 'Impact'
    temp_text.font.size = Pt(64)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.CENTER
    temp_gm.text_frame.auto_size = True



    TEXT += 1
    temp_gm = temp_page.shapes.add_textbox(Inches(5) , Inches(5) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]

    temp_text.text = '''
    This analysis examines the popularity of producer countries in Netflix content, highlighting how productions from different 
    countries impact the platform's global audience. Viewing trends and the impact of producer countries on user preferences 
    will be analyzed.
    '''
    temp_text.font.name = 'Calibri'
    temp_text.font.size = Pt(18)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.CENTER
    temp_gm.text_frame.auto_size = True



    index_list.append(PAGE)



    # DIAPOSITIVA MAPA DE CALOR COMPLETO
    # FULL HITMAP SLIDE
    PAGE += 1
    temp_page = 'page' + str(PAGE)

    TEXT = 0
    temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
    temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


    temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_7.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
    COUNT(temp_page)

    temp_page.shapes.add_picture(FOLDER_GLOBAL + 'Global_Graphs_countries_map.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)



    TEXT += 1
    temp_gm = temp_page.shapes.add_textbox(Inches(0) , Inches(0) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]

    temp_text.text = f'''
    This heatmap shows the popularity of the Producing Countries based on their total views.
    There is a small error margin due to titles with missing countries data.
    (This margin varies depending on the analyzed Data Dump.)
    Titles With No Genres: {null_titles_countries}
    Error Margin: {error_margin_countries}
    '''

    temp_text.font.name = 'times New Roman'
    temp_text.font.size = Pt(14)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.LEFT
    temp_gm.text_frame.auto_size = True



    # DIAPOSITIVA MAPA DE CALOR COMPLETO
    # FULL HITMAP SLIDE
    PAGE += 1
    temp_page = 'page' + str(PAGE)

    TEXT = 0
    temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
    temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


    temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_7.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
    COUNT(temp_page)

    temp_page.shapes.add_picture(FOLDER_GLOBAL + 'Global_Graphs_countries_map_no_top10.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)



    TEXT += 1
    temp_gm = temp_page.shapes.add_textbox(Inches(0) , Inches(0) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]

    temp_text.text = f'''
    This heatmap shows the popularity of the Producing Countries based on their total views.     IMPORTANT: In this second heatmap, we can better appreciate
    There is a small error margin due to titles with missing countries data.                                                             the popularity of lower-ranked countries
    (This margin varies depending on the analyzed Data Dump.)                                                                            by excluding the top 10 countries.
    Titles With No Genres: {null_titles_countries}
    Error Margin: {error_margin_countries}
    '''

    temp_text.font.name = 'times New Roman'
    temp_text.font.size = Pt(14)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.LEFT
    temp_gm.text_frame.auto_size = True



    # CONTEO DE LA CANTIDAD DE GRAFICOS GENERADOS PARA LOS PAISES.
    # COUNT OF THE NUMBER OF GRAPHS GENERATED FOR THE COUNTRIES.
    graphs_cuantiti_count = 1

    while True:

        if (os.path.exists(FOLDER_GLOBAL + 'Global_Graphs_countries_bars' + str(graphs_cuantiti_count) + '.png')):
            graphs_cuantiti_count += 1

        else:
            break



    # BUCLE QUE CREA TODAS LAS DIAPOSITIVAS DEL ANALISIS DE LOS PAISES SEGUN LA CANTIDAD QUE SEA NECESARIA.
    # LOOP THAT CREATES ALL THE SLIDES FOR THE COUNTRY ANALYSIS ACCORDING TO THE REQUIRED QUANTITY
    for num in range(graphs_cuantiti_count - 1):
        num = num + 1

        PAGE += 1
        temp_page = 'page' + str(PAGE)

        TEXT = 0
        temp_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
        temp_text = 'text' + str(TEXT) + '_page' + str(PAGE)


        temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
        temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_7.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
        COUNT(temp_page)



        TEXT += 1
        temp_gm = temp_page.shapes.add_textbox(Inches(0) , Inches(0) , Inches(6) , Inches(1))
        temp_text = temp_gm.text_frame.paragraphs[0]

        temp_text.text = f'''
    This bar and pie charts shows the popularity of the Producing Countries based on their total views.
    Keep in mind that the largest chart on one page is smaller than the last one on the previous page.
    There is a small error margin due to titles with missing countries data.
    (This margin varies depending on the analyzed Data Dump.)
    Titles With No Genres: {null_titles_countries}
    Error Margin: {error_margin_countries}
        '''

        temp_text.font.name = 'times New Roman'
        temp_text.font.size = Pt(14)
        temp_text.font.color.rgb = RGBColor(211, 211, 211)
        temp_text.alignment = PP_ALIGN.LEFT
        temp_gm.text_frame.auto_size = True


        # SE MIDEN EL TAMAÑO DE LA IMAGEN REAL EN PIXELESM LUEGO SE HACE UNA FORMULA PARA PASAR ESE TAMAÑO A POLGADAS Y SE APLICAN LOS DIAGRAMAS DE PIE CON LA PROPORCION CORRECTA.
        # THE ACTUAL IMAGE SIZE IS MEASURED IN PIXELS, THEN A FORMULA IS USED TO CONVERT THAT SIZE TO INCHES, AND PIE CHARTS ARE APPLIED WITH THE CORRECT PROPORTION.

        with Image.open(FOLDER_GLOBAL + 'Global_Graphs_countries_pies' + str(num) + '.png') as IMG:
            pixel_width , pixel_height = IMG.size

        inches_original_w = pixel_width / 96
        inches_original_h = pixel_height / 96

        max_w = 6
        max_h = 3

        ratio = min(max_w / inches_original_w , max_h / inches_original_h)

        correct_width = inches_original_w * ratio
        correct_height = inches_original_h * ratio



        temp_page.shapes.add_picture(FOLDER_GLOBAL + 'Global_Graphs_countries_bars' + str(num) + '.png' , left=Inches(0) , top=Inches(3) , width=diapo_.slide_width , height=(diapo_.slide_height - (diapo_.slide_height/3) ))

        temp_page.shapes.add_picture(FOLDER_GLOBAL + 'Global_Graphs_countries_pies' + str(num) + '.png' , left= Inches(9.8) , top= Inches(0) , width=Inches(correct_width) , height=Inches(correct_height))



    # PAGINAS PARA TOP50
    # PAGES FOR TOP50

    # PAGINA TITULO PARA TOP 50
    # PAGINA TITLE FOR TOP50
    PAGE += 1
    temp_page = 'page' + str(PAGE)

    temp_gm1 = 'gm1_page' + str(PAGE)
    temp_gm2 = 'gm2_page' + str(PAGE)

    temp_text1 = 'text1_page' + str(PAGE)
    temp_text2 = 'text2_page' + str(PAGE)


    temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_8.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
    COUNT(temp_page)


    temp_gm1 = temp_page.shapes.add_textbox(Inches(5) , Inches(2) , Inches(6) , Inches(1))
    temp_text1 = temp_gm1.text_frame.paragraphs[0]


    temp_text1.text = 'Analysis of \nGlobal Content \nTop_Best50_by_Views'


    temp_text1.font.name = 'Impact'
    temp_text1.font.size = Pt(64)
    temp_text1.font.color.rgb = RGBColor(211, 211, 211)
    temp_text1.alignment = PP_ALIGN.CENTER
    temp_gm1.text_frame.auto_size = True



    temp_gm2 = temp_page.shapes.add_textbox(Inches(5) , Inches(5) , Inches(6) , Inches(1))
    temp_text2 = temp_gm2.text_frame.paragraphs[0]



    temp_text2.text = """
    This analysis highlights the 50 most-viewed titles on Netflix, offering insights into the platform's top-performing content.
    These selections reflect audience preferences and showcase the titles that have captivated viewers worldwide.
    """

    temp_text2.font.name = 'Calibri'
    temp_text2.font.size = Pt(18)
    temp_text2.font.color.rgb = RGBColor(211, 211, 211)
    temp_text2.alignment = PP_ALIGN.CENTER
    temp_gm2.text_frame.auto_size = True


    index_list.append(PAGE)



    # RECUPERECION DE LOS DATOS.
    # RECOVERY OF THE DATA.


    cur.execute('''
        SELECT
            Top_50.views,
            Top_50.titles,
            Top_50.runtime,
            Top_50.votes,
            Top_50.hours_viewed,
            Top_50.countries,
            Top_50.genres,
            Top_50.format,
            Top_50.posters
        FROM
            Top_50
    ''')

    Tops50_data = cur.fetchall()


    # DESGLOSE DE LOS DATOS.
    # BREAKDOWN OF THE DATA.
    views_temp = list()
    title_ = list()
    runtime_ = list()
    votes_ = list()
    hours_temp = list()
    countries_ = list()
    genres_ = list()
    format_ = list()
    posters_ = list() 

    for row in Tops50_data:
        views_temp.append(row[0])
        title_.append(row[1])
        runtime_.append(row[2])
        votes_.append(row[3])
        hours_temp.append(row[4])
        countries_.append(row[5])
        genres_.append(row[6])
        format_.append(row[7])
        posters_.append(row[8])

    views_ = list()
    hours_ = list()
    for v , h in zip(views_temp , hours_temp):
        h = int(h)
        vi = f'{abs(v):,}'
        ho = f'{abs(h):,}'
        views_.append(vi)
        hours_.append(ho)



    ######################## CREACION DE LA CARPETA QUE CONTENDRA LOS POSTERS 
    ######################## CREATION OF THE FOLDER THAT WILL CONTAIN THE POSTERS.
    if not os.path.exists (PATH + 'global' + '_folder'):
        os.mkdir(PATH + 'global' + '_folder')

    FOLDER_POSTERS = (PATH + 'global' + '_folder' + '//')


    # DESCARGA DE LOS PORTES
    # POSTERS DOWNLOADED
    image_name_counter = 0
    for row in posters_:
        image_name_counter +=1
        
        if os.path.exists(FOLDER_POSTERS + 'Top_50' + '_' + str(image_name_counter) + '.png'):
            continue

        try:
            image_url = 'https://image.tmdb.org/t/p/w500' + row
            file_name = FOLDER_POSTERS + 'Top_50' + '_' + str(image_name_counter) + '.png'

            image_requested = requests.get(image_url)
        except:
            print('Was Not Posible to Download One Poster Image')
            continue

        if image_requested.status_code ==200:
            with open(file_name , 'wb') as file:
                file.write(image_requested.content)
        else:
            print('Was Not Posible to Download One Poster Image')




    ########## BUCLE PARA ACOMODAR CADA TITULO
    ########## LOOP FOR PLACE EVERY TITLE

    next_page = 25
    image_name_counter = 0
    for vi , ti , ru , vo , ho , co , ge , fo , po in zip(views_ , title_ , runtime_ , votes_ , hours_ , countries_ , genres_ , format_ , range(50)):

        # FILTRO DE TITULOS DEMACIADO LARGOS.
        # FILTER OF TITLES THAT ARE TOO LONG.
        if len(ti) >= 45:
            ti = ti[ :46]



        # CREACION DE UNA NUVA PAGINA PARA LA SEGUNDA MITAD
        # CREATION OF A NEW PAGE FOR DE SECOND HALF
        if next_page >= 25:
            next_page = 0
            PAGE += 1
            temp_page = 'page' + str(PAGE)

            TEXT = 0

            temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
            temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_8.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)
            COUNT(temp_page)


            # ASIGNACION DE VARIABLES A LA UBICACION DEL TEXTO Y LAS IMAGENES.
            # ASIGNATION OF VARIABLES FOR TEXT AND IMAGE PLACE.
            img_left = Inches(0.6)
            img_top = Inches(1)
            img_width = Inches(0.7)
            img_height = Inches(1)

            txt_left = img_left + img_width + Inches(0.1)
            txt_top = img_top - Inches(0.13)
            txt_width = Inches(1)
            txt_height = img_height


            jump_line = 0
            


        # CONTADOR PARA CREAR LA SEGUNDA PAGINA. Y CONTADOR SELECTOR DE IMAGENES.
        # COUNTER TO CREATE DE SECOND PAGE. AND SELECT IMAGE COUNTER.
        next_page += 1  
        image_name_counter += 1



        try:
            # IMAGEN
            # IMAGE
            image_name = FOLDER_POSTERS + 'Top_50' + '_' + str(image_name_counter) + '.png'
            temp_page.shapes.add_picture( image_name , left= img_left , top= img_top , width = img_width , height = img_height)
        except:
            # IMAGEN NEGRA
            # BLACK IMAGE
            image_name = FOLDER_PRESENTATION + str(0) + '.png'
            temp_page.shapes.add_picture( image_name , left= img_left , top= img_top , width = img_width , height = img_height)
        


        # INFORMACION DE LA PELICULA.
        # MOVIE INFORMATION
        TEXT += 1
        loop_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
        loop_text = 'text' + str(TEXT) + '_page' + str(PAGE)
        

        loop_gm = temp_page.shapes.add_textbox(txt_left , txt_top , txt_width , txt_height)
        loop_text = loop_gm.text_frame.paragraphs[0]

        loop_text.text = '\nViews: ' + str(vi) + '\nTitle: ' + str(ti) + '\nRuntime: ' + str(ru) + '\nVotes: ' + str(vo) + '\nReproducted Hours :' + str(ho) + '\nCountrie: ' + str(co) + '\nGenres: ' + str(ge) + '\n' + str(fo)

        loop_text.font.name = 'Calibri'
        loop_text.font.size = Pt(8)
        loop_text.font.color.rgb = RGBColor(211, 211, 211)
        loop_text.alignment = PP_ALIGN.LEFT
        loop_gm.text_frame.auto_size = True


        # NUMERO DE LUGAR.
        # PLACE NUMBER.
        TEXT += 1
        loop_gm = 'gm' + str(TEXT) + '_page' + str(PAGE)
        loop_text = 'text' + str(TEXT) + '_page' + str(PAGE)
        

        loop_gm = temp_page.shapes.add_textbox(img_left , img_top , txt_width , txt_height)
        loop_text = loop_gm.text_frame.paragraphs[0]

        loop_text.text = str(po + 1)

        loop_text.font.name = 'Impact'
        loop_text.font.size = Pt(16)
        loop_text.font.color.rgb = RGBColor(255, 255, 255)
        loop_text.alignment = PP_ALIGN.LEFT
        loop_gm.text_frame.auto_size = True



        # REORGANIZACION DE LO DATOS DENTRO DE LA DIAPOSITIVA. HORISONTAL
        # REORGANIZATION OF THE DATA WITHIN THE SLIDE. HORIZONTAL
        img_left = img_left + img_width + Inches(0.1) + Inches(2) + Inches(0.1)
        txt_left = img_left + img_width + Inches(0.1)


        jump_line += 1
        if jump_line >= 5:
            # REORGANIZACION DE LO DATOS DENTRO DE LA DIAPOSITIVA. VERTICAL
            # REORGANIZATION OF THE DATA WITHIN THE SLIDE. VERTICAL
            img_left = Inches(0.6)
            txt_left = img_left + img_width + Inches(0.1)

            img_top = img_top + img_height + Inches(0.5)
            txt_top = img_top - Inches(0.13)

            jump_line = 0


    conn.close()




    ############### INFORMACION DEL INDICE
    ############### INDEX INFORMATION
    conn_tv = sqlite3.connect(PATH + 'Analysis_Tv.sqlite')
    cur_tv = conn_tv.cursor()

    conn_movie = sqlite3.connect(PATH + 'Analysis_Movie.sqlite')
    cur_movie = conn_movie.cursor()


    cur_tv.execute('''
        SELECT
            Totals.views,
            Totals.hours_viewed,
            Totals.average_votes
        FROM
            Totals
    ''')
    totals = cur_tv.fetchall()
    views_tv , hours_tv , votes_tv = totals[0]

    views_tv = f'{abs(views_tv):,}'
    hours_tv = f'{abs(hours_tv):,}'
    votes_tv = str(votes_tv)[ :4]

    conn_tv.close()



    cur_movie.execute('''
        SELECT
            Totals.views,
            Totals.hours_viewed,
            Totals.average_votes
        FROM
            Totals
    ''')
    totals = cur_movie.fetchall()
    views_movie , hours_movie , votes_movie = totals[0]

    views_movie = f'{abs(views_movie):,}'
    hours_movie = f'{abs(hours_movie):,}'
    votes_movie = str(votes_movie)[ :3]

    conn_movie.close()



    temp_gm = page_4.shapes.add_textbox(Inches(0) , Inches(0.6) , Inches(6) , Inches(1))
    temp_text = temp_gm.text_frame.paragraphs[0]

    temp_text.text = f'''
    Analysis 
    of Tv/Series.
    Total Views: {views_tv}   Total Time Played: {hours_tv}   Average Votes: {votes_tv}

    Top50 Most Viewed.-----------------------------------------Page: {index_list[0]}
    Top50 Most Watched by Time Played.-------------------------Page: {index_list[1]}
    Top50 Less Viewed.-----------------------------------------Page: {index_list[2]}
    Top50 Less Watched by Time Played.-------------------------Page: {index_list[3]}
    Genres Popularity.-----------------------------------------Page: {index_list[4]}
    Producing Countries Popularity.----------------------------Page: {index_list[5]}
    Relationship Between Duration and Playback Populariry.-----Page: {index_list[6]}


    Analysis 
    of Movies.
    Total Views: {views_movie}   Total Time Played: {hours_movie}   Average Votes: {votes_movie}

    Top50 Most Viewed.-----------------------------------------Page: {index_list[7]}
    Top50 Most Watched by Time Played.-------------------------Page: {index_list[8]}
    Top50 Less Viewed.-----------------------------------------Page: {index_list[9]}
    Top50 Less Watched by Time Played.-------------------------Page: {index_list[10]}
    Genres Popularity.-----------------------------------------Page: {index_list[11]}
    Producing Countries Popularity.----------------------------Page: {index_list[12]}
    Relationship Between Duration and Playback Populariry.-----Page: {index_list[13]}


    Analysis
    of Global content.

    Analysis of Global Content.--------------------------------Page: {index_list[14]}
    Genres Popularity.-----------------------------------------Page: {index_list[15]}
    Producing Countries Popularity.----------------------------Page: {index_list[16]}
    Top50 Most Viewed.-----------------------------------------Page: {index_list[17]}
    '''

    temp_text.font.name = 'times New Roman'
    temp_text.font.size = Pt(14)
    temp_text.font.color.rgb = RGBColor(211, 211, 211)
    temp_text.alignment = PP_ALIGN.LEFT
    temp_gm.text_frame.auto_size = True



    # FIRMA
    # SING
    temp_page = diapo_.slides.add_slide(diapo_.slide_layouts[6])
    temp_page.shapes.add_picture(FOLDER_PRESENTATION + 'Pre_8.png' , 0 , 0 , width=diapo_.slide_width , height=diapo_.slide_height)


    temp_gm1 = temp_page.shapes.add_textbox(Inches(9) , Inches(5.5) , Inches(6) , Inches(1))
    temp_text1 = temp_gm1.text_frame.paragraphs[0]


    temp_text1.text = '''
    __________________________________________________

    Code Author: Daniel Sanchez Velasquez - TakuSan.
    Email: daniel.sanchez.velasquez090@gmail.com
    __________________________________________________
    '''


    temp_text1.font.name = 'Arial Narrow'
    temp_text1.font.size = Pt(18)
    temp_text1.font.color.rgb = RGBColor(211, 211, 211)
    temp_text1.alignment = PP_ALIGN.CENTER
    temp_gm1.text_frame.auto_size = True



    # GUARDADO.
    # SAVE.
    SAVE_PATH = PATH[ : -14]
    diapo_.save(SAVE_PATH + 'DATA_presentation.pptx')


    
    ############### FIN
    ############### FINISH
    print('============================== SUCCESSFUL ==============================\n')
    print('================== THE PROCESS IS COMPLETELY FINISHED ==================\n')
    print('================== THE PROCESS IS COMPLETELY FINISHED ==================\n')
    print('============================== SUCCESSFUL ==============================\n')
    print('========================================================================\n')
    print('Code Author: Daniel Sanchez Velasquez - TakuSan.')
    print('Email: daniel.sanchez.velasquez090@gmail.com' , '\n')
    print('========================================================================\n \n')
    print('The program has fully completed the normalization, analysis, and visualization.\n')
    print('You can find the presentation inside the assigned folder.\n')
    print('========================================================================\n')
    print('============================== SUCCESSFUL ==============================\n')
    print('================== THE PROCESS IS COMPLETELY FINISHED ==================\n')
    print('================== THE PROCESS IS COMPLETELY FINISHED ==================\n')
    print('============================== SUCCESSFUL ==============================\n')



    os.system(SAVE_PATH + 'DATA_presentation.pptx')


    
    return